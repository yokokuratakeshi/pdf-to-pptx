"""
PDF → PowerPoint 変換ツール
・画像モード: 各ページを高品質な画像としてスライドに変換
・編集モード: テキスト部分だけを背景から消去し、編集可能なテキストボックスを配置
"""

import streamlit as st
import fitz  # pymupdf
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io
import os

try:
    import numpy as np
    NUMPY_AVAILABLE = True
except ImportError:
    NUMPY_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image as PILImage
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# ========== ページ設定 ==========
st.set_page_config(
    page_title="PDF → PowerPoint 変換",
    page_icon="📄",
    layout="centered",
    initial_sidebar_state="collapsed",
)

# ========== カスタムCSS ==========
st.markdown("""
<style>
    .main-header { text-align: center; padding: 1.5rem 0; }
    .main-header h1 { font-size: 2.1rem; font-weight: 700; color: #1a1a2e; margin-bottom: 0.2rem; }
    .main-header p  { font-size: 0.95rem; color: #666; }
    .info-box {
        background: #f0f4ff; border-left: 4px solid #4a6cf7;
        border-radius: 0 8px 8px 0; padding: 0.9rem 1.1rem;
        margin: 0.8rem 0; font-size: 0.88rem; color: #333;
    }
    .success-box {
        background: #f0fff4; border-left: 4px solid #38a169;
        border-radius: 0 8px 8px 0; padding: 0.9rem 1.1rem; margin: 0.8rem 0;
    }
    .mode-explain {
        background: #fff8e1; border-left: 4px solid #f59e0b;
        border-radius: 0 8px 8px 0; padding: 0.9rem 1.1rem;
        margin: 0.8rem 0; font-size: 0.88rem; color: #444;
    }
    .stButton > button { width: 100%; height: 3rem; font-size: 1rem; font-weight: 600; border-radius: 8px; }
</style>
""", unsafe_allow_html=True)

# ========== ヘッダー ==========
st.markdown("""
<div class="main-header">
    <h1>📄 PDF → PowerPoint 変換</h1>
    <p>PDFをPowerPointスライドに変換します</p>
</div>
""", unsafe_allow_html=True)
st.divider()


# ========== 変換モード選択 ==========
st.markdown("### ⚙️ 変換モードを選択")
mode = st.radio(
    label="変換モード",
    options=["📷  画像モード（高品質・忠実な再現）", "✏️  編集モード（テキスト編集が可能）"],
    label_visibility="collapsed",
    horizontal=False,
)
is_edit_mode = "編集モード" in mode

if not is_edit_mode:
    st.markdown("""
    <div class="mode-explain">
        📷 <strong>画像モード</strong>：各PDFページを高品質な画像としてスライドに配置します。<br>
        見た目は原本と同じですが、テキストの直接編集はできません。
    </div>
    """, unsafe_allow_html=True)
else:
    st.markdown("""
    <div class="mode-explain">
        ✏️ <strong>編集モード</strong>：背景・図をそのまま保持しつつ、テキスト部分だけを消去して編集可能なテキストボックスに置き換えます。<br>
        &nbsp;&nbsp;• <strong>背景・図・写真・装飾</strong> → そのまま画像として表示（テキスト部分は背景色で塗りつぶし）<br>
        &nbsp;&nbsp;• <strong>テキスト</strong> → 透明な編集可能テキストボックスとして配置（クリックで編集可能）
    </div>
    """, unsafe_allow_html=True)

st.divider()

# ========== 詳細設定 ==========
st.markdown("### 🔧 設定")
col1, col2 = st.columns(2)
with col1:
    dpi = st.select_slider(
        "🎨 画像品質（DPI）",
        options=[72, 96, 144, 200, 300],
        value=144,
        help="数値が大きいほど高品質ですが、ファイルサイズも大きくなります。"
    )
with col2:
    if not is_edit_mode:
        img_fmt_label = st.selectbox(
            "🖼️ 画像形式",
            ["PNG（高品質）", "JPEG（軽量）"],
            index=0,
        )
        img_fmt = "png" if "PNG" in img_fmt_label else "jpeg"
    else:
        st.info("📌 編集モードでは背景画像は常に保持されます", icon="ℹ️")

st.divider()


# ========== ユーティリティ関数 ==========

def color_int_to_rgb(color_int):
    """PyMuPDF の色整数（0xRRGGBB）を RGBColor に変換"""
    r = (color_int >> 16) & 0xFF
    g = (color_int >> 8) & 0xFF
    b = color_int & 0xFF
    return RGBColor(r, g, b)


def clean_font_name(raw_name):
    """PDF フォント名からサブセットプレフィックスを除去して返す"""
    if not raw_name:
        return "Calibri"
    # "ABCDEF+FontName" → "FontName"
    if "+" in raw_name:
        raw_name = raw_name.split("+", 1)[1]
    # カンマ以降を除去（"Arial,Bold" → "Arial"）
    raw_name = raw_name.split(",")[0]
    # ハイフン以降を除去（"TimesNewRoman-Italic" → "TimesNewRoman"）
    raw_name = raw_name.split("-")[0]
    return raw_name.strip() or "Calibri"


def render_page_image(page, dpi, img_fmt="png"):
    """ページを画像にレンダリングして BytesIO で返す"""
    zoom = dpi / 72
    mat = fitz.Matrix(zoom, zoom)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    buf = io.BytesIO()
    buf.write(pix.tobytes(img_fmt))
    buf.seek(0)
    return buf


def render_page_without_text(page, dpi, block_bboxes):
    """
    ページをレンダリングし、テキスト領域を周囲の背景色で塗りつぶして返す。
    これにより、背景・図はそのまま残り、テキスト部分だけが消去される。

    Args:
        page: PyMuPDF のページオブジェクト
        dpi: レンダリング解像度
        block_bboxes: テキストブロックのbbox リスト（PDF座標系, ポイント単位）

    戻り値: BytesIO (PNG)
    """
    zoom = dpi / 72
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)

    if NUMPY_AVAILABLE:
        # numpy が使える場合は、テキスト領域を背景色で塗りつぶす
        img_arr = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, 3).copy()
        h, w = img_arr.shape[:2]

        for bbox in block_bboxes:
            bx0, by0, bx1, by1 = bbox
            # PDF座標 → ピクセル座標（パディングを少し加える）
            pad = max(1, int(1.5 * zoom))
            px0 = max(0, int(bx0 * zoom) - pad)
            py0 = max(0, int(by0 * zoom) - pad)
            px1 = min(w, int(bx1 * zoom) + pad)
            py1 = min(h, int(by1 * zoom) + pad)

            if px1 <= px0 or py1 <= py0:
                continue

            # 周囲のピクセルをサンプリングして背景色を推定
            sample_size = max(3, int(4 * zoom))
            edge_pixels = []

            # 上端の帯
            y_top_start = max(0, py0 - sample_size)
            y_top_end   = py0
            if y_top_start < y_top_end:
                strip = img_arr[y_top_start:y_top_end, px0:px1]
                if strip.size > 0:
                    edge_pixels.append(strip.reshape(-1, 3))

            # 下端の帯
            y_bot_start = py1
            y_bot_end   = min(h, py1 + sample_size)
            if y_bot_start < y_bot_end:
                strip = img_arr[y_bot_start:y_bot_end, px0:px1]
                if strip.size > 0:
                    edge_pixels.append(strip.reshape(-1, 3))

            # 左端の帯
            x_left_start = max(0, px0 - sample_size)
            x_left_end   = px0
            if x_left_start < x_left_end:
                strip = img_arr[py0:py1, x_left_start:x_left_end]
                if strip.size > 0:
                    edge_pixels.append(strip.reshape(-1, 3))

            # 右端の帯
            x_right_start = px1
            x_right_end   = min(w, px1 + sample_size)
            if x_right_start < x_right_end:
                strip = img_arr[py0:py1, x_right_start:x_right_end]
                if strip.size > 0:
                    edge_pixels.append(strip.reshape(-1, 3))

            if edge_pixels:
                all_pixels = np.vstack(edge_pixels)
                # 中央値で背景色を推定（外れ値に強い）
                bg_color = np.median(all_pixels, axis=0).astype(np.uint8)
            else:
                # サンプリングできない場合は白で塗る
                bg_color = np.array([255, 255, 255], dtype=np.uint8)

            img_arr[py0:py1, px0:px1] = bg_color

        # PIL → BytesIO
        result_img = PILImage.fromarray(img_arr, mode="RGB")
        buf = io.BytesIO()
        result_img.save(buf, "PNG")
        buf.seek(0)
        return buf
    else:
        # numpy が使えない場合は通常レンダリングにフォールバック
        buf = io.BytesIO()
        buf.write(pix.tobytes("png"))
        buf.seek(0)
        return buf


def ocr_page_to_textboxes(slide, page, x_scale, y_scale, ocr_dpi=200, transparent_bg=False):
    """
    OCRでページからテキストを認識してテキストボックスを追加する。

    Args:
        transparent_bg: True の場合、テキストボックスの背景を透明にする

    戻り値: (追加したテキストボックス数, テキストブロックbboxリスト)
    """
    if not OCR_AVAILABLE:
        return 0, []

    zoom = ocr_dpi / 72
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    img = PILImage.open(io.BytesIO(pix.tobytes("png")))

    try:
        data = pytesseract.image_to_data(
            img, lang="jpn+eng",
            config="--psm 6",
            output_type=pytesseract.Output.DICT
        )
    except Exception:
        try:
            # jpn が入っていない場合は eng のみで試みる
            data = pytesseract.image_to_data(
                img, lang="eng",
                config="--psm 6",
                output_type=pytesseract.Output.DICT
            )
        except Exception:
            return 0, []

    # 単語を (block_num, par_num, line_num) でグループ化して1行＝1テキストボックス
    lines = {}
    for i, text in enumerate(data["text"]):
        if not text.strip():
            continue
        conf = int(data["conf"][i])
        if conf < 30:   # 信頼度が低い認識結果はスキップ
            continue
        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        if key not in lines:
            lines[key] = {
                "words": [],
                "left":   data["left"][i],
                "top":    data["top"][i],
                "right":  data["left"][i] + data["width"][i],
                "bottom": data["top"][i]  + data["height"][i],
                "height": data["height"][i],
            }
        else:
            lines[key]["right"]  = max(lines[key]["right"],  data["left"][i] + data["width"][i])
            lines[key]["bottom"] = max(lines[key]["bottom"], data["top"][i]  + data["height"][i])
            lines[key]["height"] = max(lines[key]["height"], data["height"][i])
        lines[key]["words"].append(text)

    n_added = 0
    block_bboxes = []  # PDF座標系のbbox（render_page_without_text用）

    for line_data in lines.values():
        line_text = " ".join(line_data["words"]).strip()
        if not line_text:
            continue

        # ピクセル座標 → PDF ポイント座標 → EMU
        x0 = line_data["left"]   / zoom
        y0 = line_data["top"]    / zoom
        w  = (line_data["right"] - line_data["left"]) / zoom
        h  = (line_data["bottom"] - line_data["top"]) / zoom

        left_emu = int(x0 * x_scale)
        top_emu  = int(y0 * y_scale)
        w_emu    = int(w  * x_scale)
        h_emu    = int(h  * y_scale)

        if w_emu < 10000 or h_emu < 5000:
            continue

        # PDF座標系でbboxを記録
        block_bboxes.append((x0, y0, x0 + w, y0 + h))

        # フォントサイズをOCR行高さから推定（高さの約70%がフォントサイズに相当）
        font_pt = max(6, (line_data["height"] / zoom) * 0.70)

        txBox = slide.shapes.add_textbox(
            Emu(left_emu), Emu(top_emu),
            Emu(w_emu + 100000), Emu(h_emu + 50000)
        )
        if transparent_bg:
            # 透明背景（背景画像がテキスト部分なしで表示されるため）
            txBox.fill.background()
            txBox.line.fill.background()
        else:
            txBox.fill.solid()
            txBox.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            txBox.line.fill.background()

        tf = txBox.text_frame
        tf.word_wrap     = False
        tf.margin_left   = Emu(0)
        tf.margin_right  = Emu(0)
        tf.margin_top    = Emu(0)
        tf.margin_bottom = Emu(0)

        para = tf.paragraphs[0]
        run  = para.add_run()
        run.text      = line_text
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor(0, 0, 0)

        n_added += 1

    return n_added, block_bboxes


def convert_image_mode(doc, dpi, img_fmt, progress_bar):
    """画像モード：各ページを1枚の画像としてスライドに変換"""
    emu_per_point = 914400 / 72
    first_rect = doc[0].rect
    slide_w_emu = int(first_rect.width * emu_per_point)
    slide_h_emu = int(first_rect.height * emu_per_point)

    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank_layout = prs.slide_layouts[6]
    num_pages = len(doc)

    for i, page in enumerate(doc):
        page_rect = page.rect
        page_w_emu = int(page_rect.width * emu_per_point)
        page_h_emu = int(page_rect.height * emu_per_point)
        left = (slide_w_emu - page_w_emu) // 2
        top  = (slide_h_emu - page_h_emu) // 2

        img_buf = render_page_image(page, dpi, img_fmt)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_buf,
            Emu(max(0, left)), Emu(max(0, top)),
            width=Emu(page_w_emu), height=Emu(page_h_emu)
        )
        progress_bar.progress((i + 1) / num_pages, text=f"ページ {i+1}/{num_pages} 変換中...")

    return prs


def is_background_image(bbox, page_w, page_h, threshold=0.80):
    """ページ面積の80%以上を占める画像は背景とみなす"""
    x0, y0, x1, y1 = bbox
    img_area  = max(0, x1 - x0) * max(0, y1 - y0)
    page_area = page_w * page_h
    return page_area > 0 and (img_area / page_area) > threshold


def convert_edit_mode(doc, dpi, progress_bar):
    """
    編集モード：
    1. テキストブロックのbboxを収集
    2. テキスト部分を背景色で塗りつぶした背景画像を生成
    3. 透明背景の編集可能テキストボックスを元のテキスト位置に配置

    戻り値: (Presentation, total_textboxes)
    """
    emu_per_point = 914400 / 72
    first_rect = doc[0].rect
    slide_w_emu = int(first_rect.width * emu_per_point)
    slide_h_emu = int(first_rect.height * emu_per_point)

    prs = Presentation()
    prs.slide_width = Emu(slide_w_emu)
    prs.slide_height = Emu(slide_h_emu)
    blank_layout = prs.slide_layouts[6]
    num_pages = len(doc)

    total_textboxes = 0

    for page_idx, page in enumerate(doc):
        page_rect = page.rect
        page_w = page_rect.width
        page_h = page_rect.height

        x_scale = slide_w_emu / page_w
        y_scale = slide_h_emu / page_h

        slide = prs.slides.add_slide(blank_layout)

        # ── STEP 1: テキストブロックを抽出してbboxを収集 ──────────────
        text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_MEDIABOX_CLIP)
        text_blocks = [
            b for b in text_dict.get("blocks", [])
            if b.get("type") == 0
            and "".join(
                span.get("text", "")
                for line in b.get("lines", [])
                for span in line.get("spans", [])
            ).strip()
        ]

        use_ocr = len(text_blocks) == 0

        if use_ocr:
            # OCRモード: OCRは後で実行するため、まず通常レンダリングして仮配置
            progress_bar.progress(
                (page_idx + 0.3) / num_pages,
                text=f"ページ {page_idx+1}/{num_pages}: OCRで文字認識中..."
            )
            # OCRでbboxを先取得するため一時的に実行（テキストボックス追加なし）
            ocr_bboxes = _get_ocr_bboxes(page, dpi)
            block_bboxes_for_render = ocr_bboxes
        else:
            # 通常モード: テキストブロックのbboxをPDF座標で収集
            block_bboxes_for_render = [b["bbox"] for b in text_blocks]

        # ── STEP 2: テキスト部分を消去した背景画像を生成 ────────────
        progress_bar.progress(
            (page_idx + 0.5) / num_pages,
            text=f"ページ {page_idx+1}/{num_pages}: 背景生成中..."
        )

        if block_bboxes_for_render and NUMPY_AVAILABLE and OCR_AVAILABLE:
            # テキスト領域を背景色で塗りつぶした画像
            bg_buf = render_page_without_text(page, dpi, block_bboxes_for_render)
        else:
            # フォールバック: 通常の全ページ画像
            bg_buf = render_page_image(page, dpi, "png")

        slide.shapes.add_picture(
            bg_buf, Emu(0), Emu(0),
            width=Emu(slide_w_emu), height=Emu(slide_h_emu)
        )

        # ── STEP 3: 編集可能テキストボックスを配置 ─────────────────
        if use_ocr:
            # OCRモード: スライドにテキストボックスを追加
            n_ocr, _ = ocr_page_to_textboxes(
                slide, page, x_scale, y_scale,
                transparent_bg=True
            )
            total_textboxes += n_ocr
        else:
            # 通常モード: テキストブロックから透明テキストボックスを配置
            for block in text_blocks:
                bx0, by0, bx1, by1 = block["bbox"]
                bw, bh = bx1 - bx0, by1 - by0
                if bw <= 0 or bh <= 0:
                    continue

                left_emu = int(bx0 * x_scale)
                top_emu  = int(by0 * y_scale)
                w_emu    = int(bw  * x_scale)
                h_emu    = int(bh  * y_scale)

                if w_emu < 5000 or h_emu < 5000:
                    continue

                txBox = slide.shapes.add_textbox(
                    Emu(left_emu), Emu(top_emu),
                    Emu(w_emu + 50000), Emu(h_emu + 50000)
                )
                # 透明背景（テキスト消去済みの背景画像が透けて見える）
                txBox.fill.background()
                txBox.line.fill.background()  # 枠線なし

                tf = txBox.text_frame
                tf.word_wrap     = True
                tf.auto_size     = None
                tf.margin_left   = Emu(0)
                tf.margin_right  = Emu(0)
                tf.margin_top    = Emu(0)
                tf.margin_bottom = Emu(0)

                first_para = True
                for line in block.get("lines", []):
                    if first_para:
                        para = tf.paragraphs[0]
                        first_para = False
                    else:
                        para = tf.add_paragraph()

                    for span in line.get("spans", []):
                        text = span.get("text", "")
                        if not text:
                            continue
                        run      = para.add_run()
                        run.text = text
                        font     = run.font
                        font.size = Pt(max(1, span.get("size", 11)))
                        try:
                            font.color.rgb = color_int_to_rgb(span.get("color", 0))
                        except Exception:
                            pass
                        flags        = span.get("flags", 0)
                        font.bold    = bool(flags & 16)
                        font.italic  = bool(flags & 2)
                        try:
                            font.name = clean_font_name(span.get("font", ""))
                        except Exception:
                            pass

                total_textboxes += 1

        progress_bar.progress(
            (page_idx + 1) / num_pages,
            text=f"ページ {page_idx+1}/{num_pages} 変換中..."
        )

    return prs, total_textboxes


def _get_ocr_bboxes(page, dpi):
    """
    OCRでページのテキスト領域のbboxだけを取得する（スライドへの追加はしない）。
    render_page_without_text に渡すための事前処理。

    戻り値: PDF座標系のbboxリスト
    """
    if not OCR_AVAILABLE:
        return []

    ocr_dpi = 200
    zoom = ocr_dpi / 72
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    img = PILImage.open(io.BytesIO(pix.tobytes("png")))

    try:
        data = pytesseract.image_to_data(
            img, lang="jpn+eng",
            config="--psm 6",
            output_type=pytesseract.Output.DICT
        )
    except Exception:
        try:
            data = pytesseract.image_to_data(
                img, lang="eng",
                config="--psm 6",
                output_type=pytesseract.Output.DICT
            )
        except Exception:
            return []

    lines = {}
    for i, text in enumerate(data["text"]):
        if not text.strip():
            continue
        conf = int(data["conf"][i])
        if conf < 30:
            continue
        key = (data["block_num"][i], data["par_num"][i], data["line_num"][i])
        if key not in lines:
            lines[key] = {
                "left":   data["left"][i],
                "top":    data["top"][i],
                "right":  data["left"][i] + data["width"][i],
                "bottom": data["top"][i]  + data["height"][i],
            }
        else:
            lines[key]["right"]  = max(lines[key]["right"],  data["left"][i] + data["width"][i])
            lines[key]["bottom"] = max(lines[key]["bottom"], data["top"][i]  + data["height"][i])

    bboxes = []
    for line_data in lines.values():
        # ピクセル座標 → PDF ポイント座標
        x0 = line_data["left"]   / zoom
        y0 = line_data["top"]    / zoom
        x1 = line_data["right"]  / zoom
        y1 = line_data["bottom"] / zoom
        w  = x1 - x0
        h  = y1 - y0
        if w > 1 and h > 1:
            bboxes.append((x0, y0, x1, y1))

    return bboxes


# ========== ファイルアップロード ==========
st.markdown("### 📂 PDFファイルをアップロード")
uploaded_file = st.file_uploader(
    "ここにPDFをドラッグ＆ドロップ、またはクリックして選択",
    type=["pdf"],
)

pdf_bytes = None
if uploaded_file is not None:
    pdf_bytes = uploaded_file.read()
    try:
        doc_preview = fitz.open(stream=pdf_bytes, filetype="pdf")
        num_pages   = len(doc_preview)
        first_rect  = doc_preview[0].rect
        w_mm = first_rect.width  * 25.4 / 72
        h_mm = first_rect.height * 25.4 / 72
        doc_preview.close()

        st.markdown(f"""
        <div class="info-box">
            📋 <strong>ファイル情報</strong><br>
            &nbsp;&nbsp;• ファイル名: <code>{uploaded_file.name}</code><br>
            &nbsp;&nbsp;• ページ数: <strong>{num_pages} ページ</strong><br>
            &nbsp;&nbsp;• ページサイズ: {w_mm:.0f} × {h_mm:.0f} mm（最初のページ基準）<br>
            &nbsp;&nbsp;• ファイルサイズ: {len(pdf_bytes)/1024:.0f} KB
        </div>
        """, unsafe_allow_html=True)
    except Exception as e:
        st.error(f"PDFの読み込みに失敗しました: {e}")
        pdf_bytes = None

# ========== 変換ボタン ==========
st.markdown("### 🔄 変換")
if pdf_bytes is None:
    st.info("👆 まずPDFファイルをアップロードしてください。")
else:
    btn_label = "⚡ PowerPointに変換する（画像モード）" if not is_edit_mode else "⚡ PowerPointに変換する（編集モード）"
    if st.button(btn_label, type="primary"):
        try:
            progress_bar = st.progress(0, text="変換を開始しています...")
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")

            total_textboxes = 0

            if not is_edit_mode:
                prs = convert_image_mode(doc, dpi, img_fmt, progress_bar)
            else:
                prs, total_textboxes = convert_edit_mode(doc, dpi, progress_bar)

            doc.close()

            # PPTX をバイト列として保存
            pptx_io = io.BytesIO()
            prs.save(pptx_io)
            pptx_io.seek(0)

            progress_bar.progress(1.0, text="✅ 変換完了！")

            # 完了メッセージ
            base_name   = os.path.splitext(uploaded_file.name)[0]
            output_name = f"{base_name}.pptx"
            pptx_kb     = len(pptx_io.getvalue()) / 1024
            mode_label  = "画像モード" if not is_edit_mode else "編集モード"

            st.markdown(f"""
            <div class="success-box">
                ✅ <strong>変換が完了しました！</strong><br>
                &nbsp;&nbsp;• モード: {mode_label}<br>
                &nbsp;&nbsp;• スライド数: <strong>{num_pages} 枚</strong><br>
                &nbsp;&nbsp;• ファイルサイズ: {pptx_kb:.0f} KB
            </div>
            """, unsafe_allow_html=True)

            if is_edit_mode:
                if total_textboxes == 0:
                    st.warning(
                        "⚠️ テキストを抽出できませんでした。\n\n"
                        "**画像モード**をお試しください。"
                    )
                else:
                    st.caption(
                        f"💡 テキストボックス {total_textboxes} 個を配置しました。"
                        " PowerPointでテキストをクリックすると編集できます。"
                    )

            st.download_button(
                label=f"📥 {output_name} をダウンロード",
                data=pptx_io.getvalue(),
                file_name=output_name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary"
            )

        except Exception as e:
            st.error(f"❌ 変換中にエラーが発生しました: {str(e)}")
            st.exception(e)

# ========== フッター ==========
st.divider()
st.markdown("""
<div style="text-align:center; color:#aaa; font-size:0.78rem;">
    PDF → PowerPoint 変換ツール　|　各PDFページが1枚のスライドに変換されます
</div>
""", unsafe_allow_html=True)
